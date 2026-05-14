import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches
from pptx.oxml.ns import qn

from core.ppt_generator.utils.pptx_postprocess import (
    flatten_all_groups,
    remove_full_slide_solid_backdrops,
)


PPTX_NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}


def _make_slide_xml_blank_pptx() -> Presentation:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    prs.slides.add_slide(blank_layout)
    return prs


def _shape_tree(slide):
    return getattr(slide.shapes, "_spTree")


def _background_element(slide):
    return getattr(slide.background, "_element")


def _build_group_shape_xml(*, off, ext, child_off, child_ext, child_xml: str, rot: str = "") -> str:
    rot_attr = f' rot="{rot}"' if rot else ""
    return f"""
    <p:grpSp xmlns:p="{PPTX_NS['p']}" xmlns:a="{PPTX_NS['a']}">
      <p:nvGrpSpPr>
        <p:cNvPr id="100" name="Group 100"/>
        <p:cNvGrpSpPr/>
        <p:nvPr/>
      </p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm{rot_attr}>
          <a:off x="{off[0]}" y="{off[1]}"/>
          <a:ext cx="{ext[0]}" cy="{ext[1]}"/>
          <a:chOff x="{child_off[0]}" y="{child_off[1]}"/>
          <a:chExt cx="{child_ext[0]}" cy="{child_ext[1]}"/>
        </a:xfrm>
      </p:grpSpPr>
      {child_xml}
    </p:grpSp>
    """


def _build_rect_xml(*, off, ext, fill_hex: str = "FF0000", shape_id: int = 200) -> str:
    return f"""
    <p:sp xmlns:p="{PPTX_NS['p']}" xmlns:a="{PPTX_NS['a']}">
      <p:nvSpPr>
        <p:cNvPr id="{shape_id}" name="Rectangle {shape_id}"/>
        <p:cNvSpPr/>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="{off[0]}" y="{off[1]}"/>
          <a:ext cx="{ext[0]}" cy="{ext[1]}"/>
        </a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        <a:solidFill><a:srgbClr val="{fill_hex}"/></a:solidFill>
      </p:spPr>
    </p:sp>
    """


class FlattenAllGroupsTests(unittest.TestCase):
    def test_identity_mapping_group_dissolves_and_keeps_child_coords(self):
        from lxml import etree

        prs = _make_slide_xml_blank_pptx()
        slide = prs.slides[0]
        shape_tree = _shape_tree(slide)

        child = _build_rect_xml(off=(914400, 1828800), ext=(2743200, 1371600), shape_id=201)
        grp = _build_group_shape_xml(
            off=(914400, 1828800),
            ext=(2743200, 1371600),
            child_off=(914400, 1828800),  # identity mapping (same as drawingml_converter does)
            child_ext=(2743200, 1371600),
            child_xml=child,
        )
        shape_tree.append(etree.fromstring(grp))

        with tempfile.TemporaryDirectory() as tmp_dir:
            path = Path(tmp_dir) / "demo.pptx"
            prs.save(str(path))

            flatten_all_groups(path)

            prs2 = Presentation(str(path))
            slide2 = prs2.slides[0]
            shape_tree_two = _shape_tree(slide2)

            self.assertEqual(len(shape_tree_two.findall(qn("p:grpSp"))), 0, "group should be dissolved")
            sps = shape_tree_two.findall(qn("p:sp"))
            self.assertEqual(len(sps), 1)
            off = sps[0].find(qn("p:spPr") + "/" + qn("a:xfrm") + "/" + qn("a:off"))
            ext = sps[0].find(qn("p:spPr") + "/" + qn("a:xfrm") + "/" + qn("a:ext"))
            self.assertEqual(int(off.get("x")), 914400)
            self.assertEqual(int(off.get("y")), 1828800)
            self.assertEqual(int(ext.get("cx")), 2743200)
            self.assertEqual(int(ext.get("cy")), 1371600)

    def test_non_identity_mapping_scales_child_coords(self):
        from lxml import etree

        prs = _make_slide_xml_blank_pptx()
        slide = prs.slides[0]
        shape_tree = _shape_tree(slide)

        # Group at slide (1000, 2000) with size 4000x2000 on slide,
        # but inner coord system is 0..2000 x 0..1000 → 2x scale.
        # Child at (500, 250) inner, 1000x500 inner → should land at
        # slide (1000 + 500*2, 2000 + 250*2) = (2000, 2500), size 2000x1000.
        child = _build_rect_xml(off=(500, 250), ext=(1000, 500), shape_id=202)
        grp = _build_group_shape_xml(
            off=(1000, 2000),
            ext=(4000, 2000),
            child_off=(0, 0),
            child_ext=(2000, 1000),
            child_xml=child,
        )
        shape_tree.append(etree.fromstring(grp))

        with tempfile.TemporaryDirectory() as tmp_dir:
            path = Path(tmp_dir) / "demo.pptx"
            prs.save(str(path))

            flatten_all_groups(path)

            prs2 = Presentation(str(path))
            slide2 = prs2.slides[0]
            sp = _shape_tree(slide2).findall(qn("p:sp"))[0]
            off = sp.find(qn("p:spPr") + "/" + qn("a:xfrm") + "/" + qn("a:off"))
            ext = sp.find(qn("p:spPr") + "/" + qn("a:xfrm") + "/" + qn("a:ext"))
            self.assertEqual(int(off.get("x")), 2000)
            self.assertEqual(int(off.get("y")), 2500)
            self.assertEqual(int(ext.get("cx")), 2000)
            self.assertEqual(int(ext.get("cy")), 1000)

    def test_rotated_group_is_skipped(self):
        from lxml import etree

        prs = _make_slide_xml_blank_pptx()
        slide = prs.slides[0]
        shape_tree = _shape_tree(slide)

        child = _build_rect_xml(off=(0, 0), ext=(100, 100), shape_id=203)
        grp = _build_group_shape_xml(
            off=(0, 0), ext=(100, 100), child_off=(0, 0), child_ext=(100, 100),
            child_xml=child, rot="2700000",  # 45deg in 1/60000 units
        )
        shape_tree.append(etree.fromstring(grp))

        with tempfile.TemporaryDirectory() as tmp_dir:
            path = Path(tmp_dir) / "demo.pptx"
            prs.save(str(path))

            flatten_all_groups(path)

            prs2 = Presentation(str(path))
            shape_tree_two = _shape_tree(prs2.slides[0])
            self.assertEqual(len(shape_tree_two.findall(qn("p:grpSp"))), 1, "rotated group must remain")

    def test_nested_groups_flatten_recursively(self):
        from lxml import etree

        prs = _make_slide_xml_blank_pptx()
        slide = prs.slides[0]
        shape_tree = _shape_tree(slide)

        # Inner group: identity, holds one rect at (100, 100) size 200x200.
        inner_child = _build_rect_xml(off=(100, 100), ext=(200, 200), shape_id=204)
        inner_grp = _build_group_shape_xml(
            off=(100, 100), ext=(200, 200), child_off=(100, 100), child_ext=(200, 200),
            child_xml=inner_child,
        )
        # Outer group: identity wrap of the inner group.
        outer_grp = _build_group_shape_xml(
            off=(100, 100), ext=(200, 200), child_off=(100, 100), child_ext=(200, 200),
            child_xml=inner_grp,
        )
        shape_tree.append(etree.fromstring(outer_grp))

        with tempfile.TemporaryDirectory() as tmp_dir:
            path = Path(tmp_dir) / "demo.pptx"
            prs.save(str(path))

            flatten_all_groups(path)

            prs2 = Presentation(str(path))
            shape_tree_two = _shape_tree(prs2.slides[0])
            self.assertEqual(len(shape_tree_two.findall(qn("p:grpSp"))), 0)
            sps = shape_tree_two.findall(qn("p:sp"))
            self.assertEqual(len(sps), 1)


class RemoveFullSlideSolidBackdropTests(unittest.TestCase):
    def test_full_slide_solid_rect_lifts_to_background(self):
        from lxml import etree

        prs = _make_slide_xml_blank_pptx()
        slide = prs.slides[0]
        shape_tree = _shape_tree(slide)

        # Build a slide-sized rect at (0, 0).
        backdrop = _build_rect_xml(
            off=(0, 0), ext=(prs.slide_width, prs.slide_height),
            fill_hex="3366FF", shape_id=300,
        )
        # And a non-backdrop content rect on top.
        content = _build_rect_xml(off=(100, 100), ext=(500, 500), fill_hex="00FF00", shape_id=301)
        shape_tree.append(etree.fromstring(backdrop))
        shape_tree.append(etree.fromstring(content))

        with tempfile.TemporaryDirectory() as tmp_dir:
            path = Path(tmp_dir) / "demo.pptx"
            prs.save(str(path))

            remove_full_slide_solid_backdrops(path)

            prs2 = Presentation(str(path))
            slide2 = prs2.slides[0]
            sps = _shape_tree(slide2).findall(qn("p:sp"))
            self.assertEqual(len(sps), 1, "only the content rect must remain")
            # Background fill should have been lifted to slide.background.
            bg_srgb = _background_element(slide2).find(
                f".//{qn('a:solidFill')}/{qn('a:srgbClr')}"
            )
            self.assertIsNotNone(bg_srgb)
            self.assertEqual(bg_srgb.get("val").upper(), "3366FF")


if __name__ == "__main__":
    unittest.main()
