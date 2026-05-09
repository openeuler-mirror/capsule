import os
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from core.utils.logger import logger
from core.utils.document_parser.models import ParseResult, ImageInfo
from core.utils.document_parser.config import PptxConfig


class PptxReader:
    def __init__(self, config: PptxConfig):
        self.config = config
        self.extract_images = config.extract_images
        self.output_dir = config.output_dir
        os.makedirs(self.output_dir, exist_ok=True)

    def parse(self, file_path: str) -> ParseResult:
        logger.info("PptxReader 开始解析: {}", file_path)
        real_path = os.path.realpath(file_path)
        if not os.path.exists(real_path):
            logger.warning("PptxReader 文件不存在: {}", file_path)
            return ParseResult(text="", images=[], markdown_file="")

        try:
            prs = Presentation(real_path)
        except Exception as e:
            logger.warning("PptxReader 加载演示文稿失败: {}", e)
            return ParseResult(text="", images=[], markdown_file="")

        md_lines, image_list = self._convert_to_markdown(prs)
        md_text = "\n".join(md_lines)

        name = Path(file_path).stem
        md_filename = name + ".md"
        md_path = os.path.join(self.output_dir, md_filename)
        with open(md_path, "w", encoding="utf-8") as f:
            f.write(md_text)

        images = image_list if self.extract_images else []
        logger.info(
            "PptxReader 解析完成，文本长度: {} 字符，图片数量: {}",
            len(md_text),
            len(images),
        )
        return ParseResult(text=md_text, images=images, markdown_file=md_path)

    def _convert_to_markdown(self, prs):
        md_lines = []
        image_list = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            md_lines.append(f"## Slide {slide_idx}")
            md_lines.append("")

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            md_lines.append(text)
                            md_lines.append("")

                if self.extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_info = self._extract_image(shape)
                    if img_info:
                        image_list.append(img_info)
                        img_rel = os.path.basename(img_info.path)
                        md_lines.append(f"![image](images/{img_rel})")
                        md_lines.append("")

            md_lines.append("---")
            md_lines.append("")

        return md_lines, image_list

    def _extract_image(self, shape):
        try:
            image_blob = shape.image.blob
            image_ext = shape.image.content_type.split("/")[-1]
            if image_ext not in ("png", "jpg", "jpeg", "gif", "bmp", "svg"):
                image_ext = "png"

            image_dir = os.path.join(self.output_dir, "images")
            os.makedirs(image_dir, exist_ok=True)
            image_file = os.path.join(image_dir, f"image_{uuid.uuid4().hex[:4]}.{image_ext}")
            with open(image_file, "wb") as f:
                f.write(image_blob)
            return ImageInfo(path=os.path.realpath(image_file))
        except Exception as e:
            logger.warning("PptxReader 提取图片失败: {}", e)
            return None
