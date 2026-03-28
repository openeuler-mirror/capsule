import asyncio
import base64
import hashlib
import mimetypes
import os
import platform
import re
import shutil
import subprocess
import httpx
from pathlib import Path
from urllib.parse import urlparse

try:
    from fake_useragent import UserAgent
except ImportError:  # pragma: no cover - fallback for minimal environments
    class UserAgent:
        @property
        def random(self):
            return "Mozilla/5.0"
from PyPDF2 import PdfWriter
from pptx import Presentation
from PIL import Image

from core.utils.logger import logger
from core.utils.config import app_base_dir
from core.ppt_generator.utils.browser import BrowserManager
from core.utils.image_payload import build_image_url

UA = UserAgent()
LIBREOFFICE_DIR = Path(app_base_dir) / "libreoffice"
LINUX_SYSTEM_LIBREOFFICE_CANDIDATES = ("libreoffice26.2", "libreoffice", "soffice")
MACOS_SYSTEM_LIBREOFFICE_CANDIDATES = ("libreoffice", "soffice")
WIN_SYSTEM_LIBREOFFICE_CANDIDATES = ("soffice.com", "soffice.exe", "libreoffice", "soffice")
DEFAULT_HTML_TO_PDF_CONCURRENCY = 3
DEFAULT_RENDER_READY_TIMEOUT_MS = 20000


async def get_scale_step_value(html_path):
    """
    使用 Playwright 获取在浏览器环境中 JS 循环的最终scale step值。
    """
    async with BrowserManager.get_browser_context() as browser:
        context = await browser.new_context(viewport={'width': 1280, 'height': 720}, ignore_https_errors=True)
        page = await context.new_page()

        try:
            await page.goto(f'file://{html_path}', wait_until='networkidle', timeout=60000)
            scale_ratio = await page.evaluate("() => window.final_ration")
            logger.info(f"html {os.path.basename(html_path)} ratio: {scale_ratio}")
            return scale_ratio
        finally:
            await page.close()
            await context.close()


def sanitize_filename(name: str) -> str:
    """
    清洗文件名，替换非法字符（Windows/Linux），并将空格转换为下划线。
    保留中文、字母、数字、下划线、短横线。
    """
    cleaned = re.sub(r'[\\/*?:"<>|]', "_", name)
    cleaned = re.sub(r'\s+', "_", cleaned)
    return cleaned.strip()


async def htmls_to_pptx(html_paths: list[str], save_dir: str, filename: str = "output"):
    """
    将 HTML 路径列表转换为一个 PPTX 文件。
    """

    pdf_paths = await _batch_html_to_pdf(html_paths, save_dir)
    if not pdf_paths:
        raise Exception("没有生成任何 PDF 文件，请检查 HTML 路径是否正确。")

    merged_pdf_path = os.path.join(save_dir, f"{filename}.pdf")
    _merge_pdfs(pdf_paths, merged_pdf_path)

    for path in pdf_paths:
        if os.path.exists(path):
            os.remove(path)

    logger.info(f"正在转换 PDF 到 PPTX: {merged_pdf_path}")
    max_retries = 3
    pptx_path = ""
    for attempt in range(1, max_retries + 1):
        logger.info(f"PDF to PPTX conversion attempt {attempt}/{max_retries}...")
        pptx_path = await _libreoffice_convert_pdf_to_pptx(merged_pdf_path)

        if pptx_path:
            break

        if attempt < max_retries:
            wait_time = attempt * 2
            logger.warning(f"Attempt {attempt} failed. Retrying in {wait_time}s...")
            await asyncio.sleep(wait_time)
        else:
            logger.error("All 3 attempts to convert PDF to PPTX failed.")

    return merged_pdf_path, pptx_path


async def _batch_html_to_pdf(html_file_paths: list[str], save_dir: str) -> list[str]:
    """
    并行处理 HTML 到 PDF 的转换 (使用全局 Browser 实例)
    """
    semaphore = asyncio.Semaphore(_get_html_to_pdf_concurrency())
    async with BrowserManager.get_browser_context() as browser:
        tasks = [
            _convert_single_html_to_pdf_with_semaphore(semaphore, browser, html_path, save_dir)
            for html_path in html_file_paths
            if os.path.exists(html_path)
        ]
        results = await asyncio.gather(*tasks)

    return [path for path in results if path is not None]


def _get_html_to_pdf_concurrency() -> int:
    """
    避免多个页面同时拉取 CDN 资源，导致 Tailwind 等运行时样式尚未注入就开始导出。
    """
    raw_value = os.getenv("SLIDEA_HTML_TO_PDF_CONCURRENCY", str(DEFAULT_HTML_TO_PDF_CONCURRENCY))
    try:
        return max(1, int(raw_value))
    except (TypeError, ValueError):
        logger.warning(
            f"Invalid SLIDEA_HTML_TO_PDF_CONCURRENCY={raw_value}, fallback to {DEFAULT_HTML_TO_PDF_CONCURRENCY}"
        )
        return DEFAULT_HTML_TO_PDF_CONCURRENCY


def _get_render_ready_timeout_ms() -> int:
    raw_value = os.getenv("SLIDEA_HTML_RENDER_READY_TIMEOUT_MS", str(DEFAULT_RENDER_READY_TIMEOUT_MS))
    try:
        return max(1000, int(raw_value))
    except (TypeError, ValueError):
        logger.warning(
            f"Invalid SLIDEA_HTML_RENDER_READY_TIMEOUT_MS={raw_value}, fallback to {DEFAULT_RENDER_READY_TIMEOUT_MS}"
        )
        return DEFAULT_RENDER_READY_TIMEOUT_MS


async def _convert_single_html_to_pdf_with_semaphore(
    semaphore: asyncio.Semaphore,
    browser,
    html_file_path: str,
    save_dir: str,
) -> str | None:
    async with semaphore:
        return await _convert_single_html_to_pdf(browser, html_file_path, save_dir)


async def _wait_for_page_assets_ready(page, html_file_path: str):
    """
    显式等待字体、样式表和 Tailwind 运行时产物。
    """
    timeout_ms = _get_render_ready_timeout_ms()

    await page.wait_for_load_state("load", timeout=timeout_ms)
    await page.wait_for_function(
        """
        async () => {
            if (document.readyState !== "complete") {
                return false;
            }

            if (document.fonts && document.fonts.status !== "loaded") {
                try {
                    await document.fonts.ready;
                } catch (error) {
                    console.warn("document.fonts.ready failed", error);
                }
            }

            const stylesheetLinks = Array.from(document.querySelectorAll('link[rel="stylesheet"]'));
            const stylesheetsReady = stylesheetLinks.every((link) => {
                const href = link.getAttribute("href") || "";
                if (!href || href.startsWith("data:")) {
                    return true;
                }
                return Boolean(link.sheet);
            });
            if (!stylesheetsReady) {
                return false;
            }

            const tailwindScript = document.querySelector('script[src*="cdn.tailwindcss.com"]');
            if (!tailwindScript) {
                return true;
            }

            return Array.from(document.styleSheets).some((sheet) => {
                try {
                    const ownerNode = sheet.ownerNode;
                    const rules = Array.from(sheet.cssRules || []);
                    return (
                        ownerNode &&
                        ownerNode.tagName === "STYLE" &&
                        rules.length > 20 &&
                        rules.some((rule) => rule.cssText.includes("--tw-"))
                    );
                } catch (error) {
                    return false;
                }
            });
        }
        """,
        timeout=timeout_ms,
    )

    try:
        await page.evaluate(
            """
            async () => {
                if (typeof FontAwesome !== 'undefined' && FontAwesome && FontAwesome.dom) {
                    await FontAwesome.dom.i2svg();
                }
            }
            """
        )
        await page.wait_for_function(
            "() => !document.querySelector('[data-fa-i2svg-pending]')",
            timeout=3000,
        )
    except Exception as error:
        logger.warning(f"FontAwesome render wait skipped for {html_file_path}: {error}")

    await page.wait_for_timeout(500)


async def _convert_single_html_to_pdf(browser, html_file_path: str, save_dir: str) -> str | None:
    """
    单个页面转换逻辑
    """
    pdf_file_path = os.path.splitext(html_file_path)[0] + '.pdf'
    absolute_html_path = os.path.abspath(html_file_path)
    max_attempts = 2

    context = await browser.new_context(viewport={'width': 1500, 'height': 920}, ignore_https_errors=True)
    page = await context.new_page()

    try:
        for attempt in range(1, max_attempts + 1):
            try:
                await page.goto(f'file://{absolute_html_path}', wait_until='domcontentloaded', timeout=30000)
                await _wait_for_page_assets_ready(page, absolute_html_path)

                # 打印 PDF
                await page.pdf(
                    path=pdf_file_path,
                    width='1281px',
                    height='721px',
                    margin={'top': '0', 'right': '0', 'bottom': '0', 'left': '0'},
                    print_background=True,
                )
                logger.info(f"Successfully converted {html_file_path} -> {pdf_file_path}")
                return pdf_file_path
            except Exception as e:
                if attempt == max_attempts:
                    raise
                logger.warning(
                    f"Convert attempt {attempt}/{max_attempts} failed for {html_file_path}, retrying: {e}"
                )
                await page.wait_for_timeout(attempt * 1000)

    except Exception as e:
        logger.error(f"转换失败 {html_file_path}: {e}")
        return None

    finally:
        try:
            await page.close()
            await context.close()
        except Exception as e:
            logger.warning(f"Error closing page/context: {e}")


def _merge_pdfs(pdf_paths: list[str], output_path: str):
    """合并 PDF"""
    merger = PdfWriter()
    for pdf_path in pdf_paths:
        merger.append(pdf_path)

    with open(output_path, "wb") as f_out:
        merger.write(f_out)
    merger.close()


def _get_local_libreoffice_executable() -> Path:
    os_type = platform.system()

    if os_type == "Windows":
        program_files_dir = Path(os.environ.get("ProgramFiles", r"C:\Program Files"))
        return program_files_dir / "LibreOffice" / "program" / "soffice.com"
    if os_type == "Darwin":
        return LIBREOFFICE_DIR / "LibreOffice.app" / "Contents" / "MacOS" / "soffice"
    if os_type == "Linux":
        return LIBREOFFICE_DIR / "libreoffice-app" / "AppRun"

    raise RuntimeError(f"Unsupported operating system for LibreOffice: {os_type}")


def _get_system_libreoffice_executable() -> Path | None:
    os_type = platform.system()

    candidates = []
    if os_type == "Linux":
        candidates = LINUX_SYSTEM_LIBREOFFICE_CANDIDATES
    elif os_type == "Darwin":
        candidates = MACOS_SYSTEM_LIBREOFFICE_CANDIDATES
    elif os_type == "Windows":
        candidates = WIN_SYSTEM_LIBREOFFICE_CANDIDATES

    for candidate in candidates:
        executable = shutil.which(candidate)
        logger.info(f"Found LibreOffice {candidate} executable: {executable}")
        if executable:
            return Path(executable)
    return None


def _get_available_libreoffice_executable() -> Path | None:
    system_executable = _get_system_libreoffice_executable()
    if system_executable is not None:
        return system_executable

    local_executable = _get_local_libreoffice_executable()
    if local_executable.exists():
        return local_executable

    return None


def _build_libreoffice_pdf_to_pptx_command(
    executable: Path, file_path: str, output_dir: str
) -> list[str]:
    return [
        str(executable),
        "--headless",
        "--nologo",
        "--nolockcheck",
        "--nodefault",
        "--infilter=impress_pdf_import",
        "--convert-to",
        "pptx:Impress MS PowerPoint 2007 XML",
        "--outdir",
        output_dir,
        file_path,
    ]


async def _libreoffice_convert_pdf_to_pptx(file_path):
    """使用本地 LibreOffice 将 PDF 文件转换为 PPTX 格式。"""
    if not os.path.exists(file_path):
        logger.info(f"The file {file_path} does not exist.")
        return ""

    pptx_path = os.path.splitext(file_path)[0] + ".pptx"
    output_dir = os.path.dirname(file_path)
    executable = _get_available_libreoffice_executable()

    if executable is None or not executable.exists():
        logger.warning(
            "PDF to PPTX conversion skipped: no usable LibreOffice executable was found in the bundled directory or system PATH"
        )
        return ""

    try:
        if os.path.exists(pptx_path):
            os.remove(pptx_path)

        command = _build_libreoffice_pdf_to_pptx_command(executable, file_path, output_dir)
        logger.info(f"Client: Running local LibreOffice conversion: {' '.join(command)}")

        process = await asyncio.create_subprocess_exec(
            *command,
            cwd=str(output_dir),
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        stdout, stderr = await process.communicate()

        if process.returncode != 0:
            stdout_text = stdout.decode("utf-8", errors="ignore").strip()
            stderr_text = stderr.decode("utf-8", errors="ignore").strip()
            logger.error(
                "Local LibreOffice conversion failed: "
                f"returncode={process.returncode}, stdout={stdout_text}, stderr={stderr_text}"
            )
            return ""

        if not os.path.exists(pptx_path):
            logger.error("Local LibreOffice conversion finished but no PPTX output was generated.")
            return ""

        _force_font(pptx_path)
        _remove_bottom_layers(pptx_path)

        logger.info(f"Client: Successfully converted and saved to '{pptx_path}'")
        return pptx_path

    except Exception as e:
        logger.info(f"An error occurred while converting PDF to PPTX: {str(e)}")
        return ""


def _force_font(pptx_path, font_name="Microsoft YaHei"):
    """force all fonts in ppt transfer into Microsoft YaHei"""
    logger.info(f"Client: Forcing all fonts in {pptx_path} to {font_name}")
    prs = Presentation(pptx_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.name_far_east = font_name

    prs.save(pptx_path)
    logger.info(f"Client: Successfully updated font in '{pptx_path}'")


def _remove_bottom_layers(pptx_path):
    """
    Remove the bottom two shapes (Z-order wise) from every slide in the pptx.
    python-pptx shapes[0] is the bottom-most layer.
    """
    logger.info(f"Client: Removing bottom 2 shapes from slides in {pptx_path}")
    prs = Presentation(pptx_path)

    for _, slide in enumerate(prs.slides):
        for _ in range(2):
            if len(slide.shapes) > 0:
                shape_to_delete = slide.shapes[0]
                sp = shape_to_delete._element
                sp.getparent().remove(sp)
            else:
                break

    prs.save(pptx_path)
    logger.info(f"Client: Successfully removed bottom shapes in '{pptx_path}'")


def _extract_web_image_description(image: dict, image_query: str) -> str:
    for key in ["description", "image_description", "content", "caption", "alt", "title"]:
        value = image.get(key)
        if isinstance(value, str) and value.strip():
            return value.strip()
    return image_query


async def get_web_images_content(image_query_list, image_list, save_dir):
    """download image of images query"""
    result = []
    img_list = []
    image_descriptions = {}
    images_dir = os.path.join(save_dir, "images")
    if not os.path.exists(images_dir):
        os.makedirs(images_dir)

    nested_results = await _execute_download_images(image_list, images_dir)

    for download_result, image_result in zip(nested_results, image_list):
        for download_img, image in zip(download_result, image_result):
            image['url'] = download_img

    for image_query, image_paths in zip(image_query_list, image_list):
        image_paths = [item for item in image_paths if item['url'] is not None]
        result.append(f"图片'{image_query}'的下载结果：{image_paths}")
        for image in image_paths:
            image_path = os.path.join(save_dir, image['url'])
            img_list.append(image_path)
            image_descriptions[image_path] = _extract_web_image_description(image, image_query)

    return "\n".join(result), img_list, image_descriptions


async def _execute_download_images(image_list, images_dir):
    """execute download images"""
    tasks = []
    group_sizes = []
    for image_result in image_list:
        group_sizes.append(len(image_result))
        for image in image_result:
            img_url = image['url']
            task = asyncio.create_task(download_image(img_url, images_dir))
            tasks.append(task)
    flat_results = await asyncio.gather(*tasks)
    nested_results = []
    current_pos = 0
    for size in group_sizes:
        chunk = flat_results[current_pos: current_pos + size]
        nested_results.append(chunk)
        current_pos += size
    return nested_results


def _ensure_placeholder_image(image_dir: str) -> str:
    os.makedirs(image_dir, exist_ok=True)
    placeholder_path = os.path.join(image_dir, "placeholder.png")
    if os.path.exists(placeholder_path):
        return placeholder_path

    # 1x1 transparent PNG
    placeholder_b64 = (
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR4nGNgYAAAAAMAASsJTYQAAAAASUVORK5CYII="
    )
    with open(placeholder_path, "wb") as f:
        f.write(base64.b64decode(placeholder_b64))
    return placeholder_path


async def download_image(img_url, image_dir):
    """download image and return img used in html"""
    if img_url.startswith("//"):
        img_url = "https:" + img_url
    headers = {
        "User-Agent": UA.random,
        "Accept": "image/avif,image/webp,image/apng,image/*,*/*;q=0.8",
        "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8",
    }
    parsed_uri = urlparse(img_url)
    headers["Referer"] = f"{parsed_uri.scheme}://{parsed_uri.netloc}/"

    try:
        # 1. 先执行下载请求
        async with httpx.AsyncClient(verify=False) as client:
            response = await client.get(
                img_url,
                headers=headers,
                timeout=20.0,
                follow_redirects=True,
            )
            if response.status_code == 403:
                # 403 常见为防盗链，去掉 Referer 重试
                headers.pop("Referer", None)
                response = await client.get(
                    img_url,
                    headers=headers,
                    timeout=20.0,
                    follow_redirects=True,
                )
            response.raise_for_status()

            # 2. 从响应头中获取内容类型 (Content-Type)
            content_type = response.headers.get("Content-Type")
            if not content_type or not content_type.startswith("image/"):
                # Some hosts return application/octet-stream for images
                if not img_url.lower().endswith((
                    ".png",
                    ".jpg",
                    ".jpeg",
                    ".webp",
                    ".gif",
                    ".bmp",
                    ".svg",
                )):
                    logger.debug(
                        f"下载失败: {img_url}, Content-Type非图片格式: {content_type}"
                    )
                    return _ensure_placeholder_image(image_dir)

            # 3. 使用mimetypes库将 'image/jpeg' 转换为 '.jpg' 等扩展名
            file_ext = mimetypes.guess_extension(content_type.split(";")[0])
            if not file_ext:
                # 如果 mimetypes 无法识别，提供一个简单的备用方案
                subtype = content_type.split("/")[-1].split(";")[0]
                file_ext = f".{subtype}"
                logger.debug(
                    f"无法从 '{content_type}' 自动推断扩展名, 回退使用: '{file_ext}'"
                )

            # 常见修正: .jpe -> .jpg
            if file_ext == ".jpe":
                file_ext = ".jpg"

            # 4. 生成文件名和路径
            filename = f"{hashlib.md5(img_url.encode()).hexdigest()}{file_ext}"
            local_path = os.path.join(image_dir, filename)

            # 5. 保存到本地
            with open(local_path, "wb") as file:
                file.write(response.content)

            # 6. 转换不支持的格式（avif/webp -> jpg）
            if file_ext in (".avif", ".webp"):
                try:
                    jpg_path = os.path.splitext(local_path)[0] + ".jpg"
                    with Image.open(local_path) as im:
                        im = im.convert("RGB")
                        im.save(jpg_path, "JPEG", quality=90)
                    return jpg_path
                except Exception as e:
                    logger.warning(f"转换图片格式失败: {local_path} - {e}")

            return local_path
    except Exception as e:
        logger.debug(f"下载时发生未知错误: {img_url} - {str(e)}")
        return _ensure_placeholder_image(image_dir)
