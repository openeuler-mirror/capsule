"""Post-processing helpers shared by the HTML and SVG → PPTX export routes.

Both routes have the same problem: the source format expresses the slide
backdrop as the bottom-most full-slide solid-fill shape (LibreOffice produces
those when converting PDF→PPTX; the LLM-generated SVG files start with a
``<rect width="100%" height="100%" fill="#x">``). Either way the result is a
real shape on the slide instead of the slide's own background fill, which
hurts editability in PowerPoint.

``remove_full_slide_solid_backdrops`` walks the z-order from the bottom up,
drops every shape that is a full-slide solid fill, and lifts the last removed
color onto ``slide.background.fill`` so the page keeps its backdrop.

``flatten_all_groups`` recursively dissolves every ``p:grpSp`` on every slide,
hoisting each group's children onto the slide (or the parent group) with
absolute coordinates. Run this before backdrop detection on SVG-built PPTX so
group bounding boxes never falsely match "full-slide solid fill" and so a
backdrop rect that happens to be wrapped in a group still gets stripped.
"""

from dataclasses import dataclass

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

from core.utils.logger import logger


_DRAWABLE_TAGS = frozenset(qn(t) for t in ("p:sp", "p:cxnSp", "p:pic", "p:grpSp", "p:graphicFrame"))


@dataclass(frozen=True)
class GroupTransform:
    """Coordinate transform from a group inner canvas into its parent."""

    off_x: int
    off_y: int
    child_off_x: int
    child_off_y: int
    scale_x: float
    scale_y: float


def remove_full_slide_solid_backdrops(pptx_path) -> None:
    """Strip full-slide solid-fill backdrop shapes in-place and lift to slide background.

    python-pptx ``shapes[0]`` is the bottom-most layer. We pop shapes from the
    bottom while they are full-slide solid fills; the last removed color is
    promoted to the slide's background fill.
    """
    logger.info(f"Stripping full-slide solid backdrops in {pptx_path}")
    prs = Presentation(str(pptx_path))
    slide_w, slide_h = prs.slide_width, prs.slide_height

    for slide in prs.slides:
        last_rgb = None
        removed = 0
        while len(slide.shapes) > 0:
            shape = slide.shapes[0]
            rgb = _full_slide_solid_fill_rgb(shape, slide_w, slide_h)
            if rgb is None:
                break
            last_rgb = rgb
            sp = getattr(shape, "_element")
            sp.getparent().remove(sp)
            removed += 1

        if last_rgb is not None:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = last_rgb
        logger.debug(f"removed {removed} backdrop shape(s) from slide")

    prs.save(str(pptx_path))
    logger.info(f"Successfully cleaned bottom backdrop shapes in '{pptx_path}'")


def _full_slide_solid_fill_rgb(shape, slide_w, slide_h, tol_ratio=0.1):
    """Return shape's solid-fill RGB iff it has one AND covers the whole slide
    within ``tol_ratio`` of slide dimensions. Otherwise None."""
    rgb = _extract_solid_fill_rgb(shape)
    if rgb is None:
        return None
    try:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
    except AttributeError:
        return None
    if None in (left, top, width, height):
        return None
    w_tol = slide_w * tol_ratio
    h_tol = slide_h * tol_ratio
    covers_slide = (
        left <= w_tol
        and top <= h_tol
        and left + width >= slide_w - w_tol
        and top + height >= slide_h - h_tol
    )
    return rgb if covers_slide else None


def _extract_solid_fill_rgb(shape):
    """Return the shape's solid-fill RGB color, or None if absent/unsupported."""
    shape_element = getattr(shape, "_element")
    srgb = shape_element.find(f".//{qn('a:solidFill')}/{qn('a:srgbClr')}")
    if srgb is None:
        return None
    val = srgb.get("val")
    if not val:
        return None
    try:
        return RGBColor.from_string(val)
    except ValueError:
        return None


def flatten_all_groups(pptx_path) -> None:
    """Recursively dissolve every ``p:grpSp`` on every slide.

    Each group's drawable children are spliced into the group's parent in the
    same z-order, with their offsets/extents transformed from the group's
    inner coordinate system into the parent's coordinate system. The group's
    rotation is honored only when zero — groups with non-zero rotation are
    left intact (rotating each child around the group's pivot is geometrically
    non-trivial and rare in LLM-generated SVG).

    Idempotent: running on a PPTX with no groups does nothing.
    """
    logger.info(f"Flattening grouped shapes in {pptx_path}")
    prs = Presentation(str(pptx_path))
    total = 0
    for slide in prs.slides:
        shape_tree = getattr(slide.shapes, "_spTree")
        # Each pass dissolves immediate-child groups in shape_tree. Children of a
        # dissolved group that are themselves groups become new immediate
        # children, so loop until none remain.
        while True:
            n = _ungroup_one_pass(shape_tree)
            total += n
            if n == 0:
                break
    prs.save(str(pptx_path))
    logger.info(f"Flattened {total} group(s) in '{pptx_path}'")


def _ungroup_one_pass(parent) -> int:
    """Dissolve every immediate-child ``p:grpSp`` of ``parent``. Returns count."""
    group_shape_tag = qn("p:grpSp")
    flattened = 0
    for group in [c for c in list(parent) if c.tag == group_shape_tag]:
        if _flatten_group(group):
            flattened += 1
    return flattened


def _flatten_group(group) -> bool:
    """Replace ``group`` with its drawable children in ``group``'s parent.

    Returns True iff the group was successfully dissolved. False if the group
    has rotation (skipped to avoid wrong child positions) or if the group has
    no parent (already detached).
    """
    parent = group.getparent()
    if parent is None:
        return False

    group_xfrm = group.find(qn("p:grpSpPr") + "/" + qn("a:xfrm"))
    rot = group_xfrm.get("rot") if group_xfrm is not None else None
    if rot and int(rot) != 0:
        logger.debug("Skip flattening rotated group (rot=%s)", rot)
        return False

    if group_xfrm is None:
        transform = GroupTransform(0, 0, 0, 0, 1.0, 1.0)
    else:
        off = group_xfrm.find(qn("a:off"))
        ext = group_xfrm.find(qn("a:ext"))
        child_off = group_xfrm.find(qn("a:chOff"))
        child_ext = group_xfrm.find(qn("a:chExt"))
        off_x = int(off.get("x", "0")) if off is not None else 0
        off_y = int(off.get("y", "0")) if off is not None else 0
        ex = int(ext.get("cx", "0")) if ext is not None else 0
        ey = int(ext.get("cy", "0")) if ext is not None else 0
        child_off_x = int(child_off.get("x", "0")) if child_off is not None else 0
        child_off_y = int(child_off.get("y", "0")) if child_off is not None else 0
        child_ext_x = int(child_ext.get("cx", str(ex or 1))) if child_ext is not None else (ex or 1)
        child_ext_y = int(child_ext.get("cy", str(ey or 1))) if child_ext is not None else (ey or 1)
        transform = GroupTransform(
            off_x,
            off_y,
            child_off_x,
            child_off_y,
            (ex / child_ext_x) if child_ext_x else 1.0,
            (ey / child_ext_y) if child_ext_y else 1.0,
        )

    children = [c for c in list(group) if c.tag in _DRAWABLE_TAGS]
    insert_idx = list(parent).index(group)
    for child in children:
        _apply_group_transform(child, transform)
        parent.insert(insert_idx, child)
        insert_idx += 1
    parent.remove(group)
    return True


def _apply_group_transform(child, transform: GroupTransform) -> None:
    """Map child's xfrm offset/extent from group's inner coords into parent coords.

    Formula: ``slide_x = group.off.x + (child.x - group.chOff.x) * (group.ext / group.chExt)``.
    For nested ``p:grpSp`` children we transform their own off/ext only — chOff /
    chExt define an internal coordinate system independent of how the nested
    group is placed in its parent and stay untouched.
    """
    xfrm = _find_xfrm(child)
    if xfrm is None:
        return
    off = xfrm.find(qn("a:off"))
    ext = xfrm.find(qn("a:ext"))
    if off is not None:
        x = int(off.get("x", "0"))
        y = int(off.get("y", "0"))
        off.set("x", str(round(transform.off_x + (x - transform.child_off_x) * transform.scale_x)))
        off.set("y", str(round(transform.off_y + (y - transform.child_off_y) * transform.scale_y)))
    if ext is not None:
        cx = int(ext.get("cx", "0"))
        cy = int(ext.get("cy", "0"))
        ext.set("cx", str(max(round(cx * transform.scale_x), 1)))
        ext.set("cy", str(max(round(cy * transform.scale_y), 1)))


def _find_xfrm(child):
    """Return the child element's xfrm element, or None if absent.

    Drawable XML location varies by tag::

        p:sp / p:cxnSp / p:pic → p:spPr/a:xfrm
        p:grpSp                → p:grpSpPr/a:xfrm
        p:graphicFrame         → p:xfrm
    """
    if child.tag in (qn("p:sp"), qn("p:cxnSp"), qn("p:pic")):
        return child.find(qn("p:spPr") + "/" + qn("a:xfrm"))
    if child.tag == qn("p:grpSp"):
        return child.find(qn("p:grpSpPr") + "/" + qn("a:xfrm"))
    if child.tag == qn("p:graphicFrame"):
        return child.find(qn("p:xfrm"))
    return None
